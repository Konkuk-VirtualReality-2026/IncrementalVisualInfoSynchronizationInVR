"""
중간 발표 PPT — 발표자 구두 설명 보조용 (키워드 중심, 텍스트 최소화)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BLUE   = RGBColor(0x16, 0x45, 0xA8)
GRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
GREEN  = RGBColor(0x1A, 0x7A, 0x1A)
ORANGE = RGBColor(0xCC, 0x44, 0x00)
PURPLE = RGBColor(0x6A, 0x0D, 0xAD)
BGBLUE = RGBColor(0xEE, 0xF4, 0xFF)
BGGRN  = RGBColor(0xEE, 0xF8, 0xEE)
BGPRP  = RGBColor(0xF5, 0xEE, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
IMG_DIR = r"D:\Git\SynchronizationInVR\Doc\중간보고 사진"

def img(name): return os.path.join(IMG_DIR, name)

# ── 기본 헬퍼 ───────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(slide, text, left, top, w, h,
        size=18, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "맑은 고딕"
    return tb

def multibox(slide, lines, left, top, w, h, size=18, color=BLACK):
    """lines: list of str or dict(text, bold, color, size)"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            d = {"text": item}
        else:
            d = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = d.get("align", PP_ALIGN.LEFT)
        r = p.add_run()
        r.text = d.get("text", "")
        r.font.size  = Pt(d.get("size", size))
        r.font.bold  = d.get("bold", False)
        r.font.color.rgb = d.get("color", color)
        r.font.name  = "맑은 고딕"
    return tb

def line(slide, top, left=0.5, right=12.83, color=LGRAY):
    rect = slide.shapes.add_shape(1, Inches(left), top, Inches(right - left), Pt(1.5))
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()

def pill(slide, text, left, top, w=Inches(2.0), h=Inches(0.45),
         bg=BLUE, fg=WHITE, size=14):
    rect = slide.shapes.add_shape(1, left, top, w, h)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    tf = rect.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = True; r.font.color.rgb = fg; r.font.name = "맑은 고딕"

def safe_pic(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        try:
            kw = {}
            if width:  kw["width"]  = width
            if height: kw["height"] = height
            return slide.shapes.add_picture(path, left, top, **kw)
        except Exception as e:
            print(f"  [WARN] {os.path.basename(path)}: {e}")
    # 대체 박스
    w = width or Inches(4); h = height or Inches(3)
    r = slide.shapes.add_shape(1, left, top, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    r.line.color.rgb = LGRAY
    tf = r.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    rn = tf.paragraphs[0].add_run()
    rn.text = os.path.basename(path); rn.font.size = Pt(10); rn.font.color.rgb = GRAY
    return r

def rect_bg(slide, left, top, w, h, color):
    r = slide.shapes.add_shape(1, left, top, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    return r

def header(slide, title, sub=None):
    box(slide, title, Inches(0.6), Inches(0.2), Inches(12), Inches(0.75),
        size=30, bold=True, color=DARK)
    if sub:
        box(slide, sub, Inches(0.6), Inches(0.88), Inches(12), Inches(0.4),
            size=15, color=GRAY, italic=True)
    line(slide, Inches(1.22) if sub else Inches(0.9))


# ════════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════════
def s_title(prs):
    s = blank(prs)
    box(s, "건국대학교  ·  2026-1  가상현실[3221]  ·  중간 발표",
        Inches(0.6), Inches(1.1), Inches(12), Inches(0.55),
        size=16, color=GRAY, align=PP_ALIGN.CENTER)
    box(s, "VR 멀미 저감을 위한\n교차 감각 기반\n점진적 시각 동기화 기법",
        Inches(0.6), Inches(1.75), Inches(12), Inches(2.6),
        size=40, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    box(s, "Progressive Visual Synchronization via Cross-modal Sensory Adaptation",
        Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
        size=15, color=BLUE, italic=True, align=PP_ALIGN.CENTER)
    line(s, Inches(5.05))
    box(s, "3팀  |  정근녕  ·  이수민  ·  곽경민  |  담당: 김형석 교수님",
        Inches(0.6), Inches(5.2), Inches(12), Inches(0.45),
        size=17, color=DARK, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 슬라이드 2: 목차
# ════════════════════════════════════════════════════════════════
def s_agenda(prs):
    s = blank(prs)
    header(s, "목  차")
    items = [
        ("01", "왜 VR이 멀미를 유발하는가?",       "연구 배경"),
        ("02", "우리의 해결 방법: 3단계 적응",       "연구 목표 및 내용"),
        ("03", "지금까지 만든 것들",                 "구현 현황"),
        ("04", "앞으로 할 것들",                     "향후 계획 및 실험 일정"),
        ("05", "누가 무엇을 했는가",                 "역할 분담"),
    ]
    for i, (num, title, sub) in enumerate(items):
        top = Inches(1.4) + i * Inches(1.05)
        pill(s, num, Inches(0.6), top + Inches(0.05),
             w=Inches(0.65), h=Inches(0.45), size=16)
        box(s, title, Inches(1.45), top, Inches(7.5), Inches(0.52),
            size=21, bold=True, color=DARK)
        box(s, sub,   Inches(9.2),  top, Inches(3.7), Inches(0.52),
            size=15, color=GRAY, align=PP_ALIGN.RIGHT)
        if i < len(items) - 1:
            line(s, top + Inches(0.7), left=0.6, right=12.73, color=LGRAY)


# ════════════════════════════════════════════════════════════════
# 슬라이드 3: 배경 — 왜 VR이 멀미를 유발하나
# ════════════════════════════════════════════════════════════════
def s_bg1(prs):
    s = blank(prs)
    header(s, "왜 VR은 멀미를 유발하는가?")

    # 중앙 큰 명제
    box(s, "뇌의 예측  ≠  실제 감각 입력",
        Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.85),
        size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    line(s, Inches(2.35), left=2.0, right=11.33, color=BLUE)

    # 좌 — 현실 세계
    rect_bg(s, Inches(0.5), Inches(2.5), Inches(5.7), Inches(3.6), BGBLUE)
    box(s, "현실 세계", Inches(0.7), Inches(2.6), Inches(5.3), Inches(0.5),
        size=18, bold=True, color=BLUE)
    multibox(s, [
        "오랜 경험으로 형성된",
        {"text": "'감각 예측 모델'", "bold": True, "color": BLUE, "size": 20},
        "",
        "시각  ↔  전정기관  ↔  촉각",
        "→ 완벽하게 일치하는 상태",
    ], Inches(0.7), Inches(3.15), Inches(5.3), Inches(2.7), size=17, color=DARK)

    # 가운데 충돌 표시
    box(s, "⚡\n충돌",
        Inches(6.3), Inches(3.1), Inches(0.9), Inches(1.2),
        size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # 우 — VR 세계
    rect_bg(s, Inches(7.1), Inches(2.5), Inches(5.7), Inches(3.6), RGBColor(0xFF,0xF2,0xEE))
    box(s, "VR 세계", Inches(7.3), Inches(2.6), Inches(5.3), Inches(0.5),
        size=18, bold=True, color=ORANGE)
    multibox(s, [
        "헤드셋 착용 즉시",
        {"text": "기존 예측 모델과 전혀 다른", "bold": True, "color": ORANGE, "size": 20},
        "",
        "시각은 이동  ↔  몸은 정지",
        "→ 멀미(Cybersickness) 발생",
    ], Inches(7.3), Inches(3.15), Inches(5.3), Inches(2.7), size=17, color=DARK)

    box(s, "기존 해결책(vignetting·rest frame·하드웨어)은 '몰입감'을 희생하는 trade-off가 있다",
        Inches(0.5), Inches(6.25), Inches(12.33), Inches(0.5),
        size=14, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 슬라이드 4: 배경 — 우리의 접근
# ════════════════════════════════════════════════════════════════
def s_bg2(prs):
    s = blank(prs)
    header(s, "우리의 접근: 뇌를 적응시킨다")

    box(s, "멀미를 '피하는' 것이 아니라\n뇌의 예측 모델 자체를 바꾼다",
        Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.4),
        size=30, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    line(s, Inches(2.95), left=1.5, right=11.83, color=PURPLE)

    # 3개 근거 카드
    cards = [
        (BGPRP, PURPLE, "뇌 가소성\n(Neuroplasticity)",
         "반복 자극으로\n뇌가 새 환경에 적응"),
        (BGBLUE, BLUE, "선행 노출\n(Pre-exposure)",
         "VR 진입 전\n감각 예측 모델 재구성"),
        (BGGRN, GREEN, "단계적 동기화\n(Progressive Sync)",
         "감각 정보를\n점진적으로 늘려간다"),
    ]
    for i, (bg, col, title, desc) in enumerate(cards):
        lf = 0.5 + i * 4.25
        left = Inches(lf)
        rect_bg(s, left, Inches(3.15), Inches(3.9), Inches(3.8), bg)
        box(s, title, left + Inches(0.2), Inches(3.3), Inches(3.5), Inches(0.9),
            size=19, bold=True, color=col)
        line(s, Inches(4.3), left=lf + 0.2, right=lf + 3.7, color=col)
        box(s, desc, left + Inches(0.2), Inches(4.45), Inches(3.5), Inches(2.1),
            size=18, color=DARK)


# ════════════════════════════════════════════════════════════════
# 슬라이드 5: 3단계 방법론 개요
# ════════════════════════════════════════════════════════════════
def s_overview(prs):
    s = blank(prs)
    header(s, "3단계 점진적 시각 동기화",
           "VR 진입 → Phase 1 (30s) → Phase 2 (30s) → Phase 3 (30s) → 본 콘텐츠")

    phases = [
        (BGPRP, PURPLE, "Phase 1",
         "완전한 어둠",
         "소리 + 진동으로\n공간을 상상하게 한다",
         img("스크린샷 2026-05-18 112418.png")),
        (BGBLUE, BLUE, "Phase 2",
         "윤곽선만",
         "상상한 공간과\n시각 정보를 연결한다",
         img("비교군1png.png")),
        (BGGRN, GREEN, "Phase 3",
         "점점 선명하게",
         "새 예측 모델 위에\n완전한 VR을 입힌다",
         img("비교군2png.png")),
    ]

    for i, (bg, col, ph, kw, desc, img_path) in enumerate(phases):
        left = Inches(0.5) + i * Inches(4.25)
        rect_bg(s, left, Inches(1.4), Inches(3.9), Inches(5.6), bg)

        pill(s, ph, left + Inches(0.18), Inches(1.55),
             w=Inches(1.2), h=Inches(0.4), bg=col, size=13)
        box(s, kw,   left + Inches(0.18), Inches(2.08),
            Inches(3.54), Inches(0.7), size=22, bold=True, color=col)
        box(s, desc, left + Inches(0.18), Inches(2.82),
            Inches(3.54), Inches(0.85), size=15, color=DARK)
        safe_pic(s, img_path,
                 left + Inches(0.1), Inches(3.75),
                 width=Inches(3.7), height=Inches(2.85))


# ════════════════════════════════════════════════════════════════
# 슬라이드 6: Phase 1 상세
# ════════════════════════════════════════════════════════════════
def s_phase1(prs):
    s = blank(prs)
    header(s, "Phase 1 — 완전한 어둠 속에서", "30초 | 시각 100% 차단")

    # 좌 — 3가지 핵심 포인트
    pill(s, "핵심 질문", Inches(0.6), Inches(1.45), w=Inches(1.8), h=Inches(0.42), bg=PURPLE)
    box(s, "어둠 속에서도\n공간을 느낄 수 있는가?",
        Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.2),
        size=26, bold=True, color=DARK)

    line(s, Inches(3.3), left=0.6, right=6.6, color=PURPLE)

    multibox(s, [
        {"text": "   소리", "bold": True, "color": PURPLE, "size": 20},
        {"text": "   앰비언트 오디오로 공간감 제공", "size": 16, "color": GRAY},
        "",
        {"text": "   진동", "bold": True, "color": PURPLE, "size": 20},
        {"text": "   컨트롤러 햅틱 — 심박수 리듬(1.5s)", "size": 16, "color": GRAY},
        "",
        {"text": "   상상", "bold": True, "color": PURPLE, "size": 20},
        {"text": "   피실험자가 스스로 공간 구조를 추론", "size": 16, "color": GRAY},
    ], Inches(0.6), Inches(3.45), Inches(6.0), Inches(3.5), size=18)

    # 우 — 스크린샷
    safe_pic(s, img("스크린샷 2026-05-18 112418.png"),
             Inches(6.9), Inches(1.38), width=Inches(6.1))
    safe_pic(s, img("스크린샷 2026-05-18 112502.png"),
             Inches(6.9), Inches(4.45), width=Inches(6.1), height=Inches(2.7))


# ════════════════════════════════════════════════════════════════
# 슬라이드 7: Phase 2 상세
# ════════════════════════════════════════════════════════════════
def s_phase2(prs):
    s = blank(prs)
    header(s, "Phase 2 — 윤곽선만 보여준다", "30초 | fidelity = 0.3")

    pill(s, "핵심 질문", Inches(0.6), Inches(1.45), w=Inches(1.8), h=Inches(0.42), bg=BLUE)
    box(s, "상상했던 공간과\n지금 보이는 윤곽선이 맞는가?",
        Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.2),
        size=26, bold=True, color=DARK)

    line(s, Inches(3.3), left=0.6, right=6.6, color=BLUE)

    multibox(s, [
        {"text": "   최소 시각", "bold": True, "color": BLUE, "size": 20},
        {"text": "   엣지·윤곽선만 표시 (표면은 검정)", "size": 16, "color": GRAY},
        "",
        {"text": "   매칭", "bold": True, "color": BLUE, "size": 20},
        {"text": "   Phase 1 공간 인식 + 시각 정보 연결", "size": 16, "color": GRAY},
        "",
        {"text": "   전환", "bold": True, "color": BLUE, "size": 20},
        {"text": "   현실 예측치 → 가상 공간 예측치 형성", "size": 16, "color": GRAY},
    ], Inches(0.6), Inches(3.45), Inches(6.0), Inches(3.5), size=18)

    safe_pic(s, img("비교군1png.png"),
             Inches(6.9), Inches(1.38), width=Inches(6.1))
    safe_pic(s, img("비교군.png"),
             Inches(6.9), Inches(4.45), width=Inches(6.1), height=Inches(2.7))


# ════════════════════════════════════════════════════════════════
# 슬라이드 8: Phase 3 상세
# ════════════════════════════════════════════════════════════════
def s_phase3(prs):
    s = blank(prs)
    header(s, "Phase 3 — 점점 선명해진다", "30초 | fidelity 0.3 → 1.0")

    pill(s, "결과", Inches(0.6), Inches(1.45), w=Inches(1.8), h=Inches(0.42), bg=GREEN)
    box(s, "뇌가 이미 가상 공간에\n익숙해진 상태에서 VR이 시작된다",
        Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.3),
        size=24, bold=True, color=DARK)

    line(s, Inches(3.4), left=0.6, right=6.6, color=GREEN)

    multibox(s, [
        {"text": "   fidelity 0.3 → 1.0", "bold": True, "color": GREEN, "size": 20},
        {"text": "   텍스처·광원·색상 30초에 걸쳐 점진 복원", "size": 16, "color": GRAY},
        "",
        {"text": "   완전한 VR", "bold": True, "color": GREEN, "size": 20},
        {"text": "   이미 형성된 예측 모델 위에 현실감을 입힘", "size": 16, "color": GRAY},
        "",
        {"text": "   → AimTrainer 시작", "bold": True, "color": DARK, "size": 18},
    ], Inches(0.6), Inches(3.55), Inches(6.0), Inches(3.3), size=18)

    safe_pic(s, img("비교군2png.png"),
             Inches(6.9), Inches(1.38), width=Inches(6.1))


# ════════════════════════════════════════════════════════════════
# 슬라이드 9: 구현 현황
# ════════════════════════════════════════════════════════════════
def s_impl(prs):
    s = blank(prs)
    header(s, "지금까지 만든 것들", "핵심 기능 구현 완료 — 실험 환경 구축 직전 단계")

    items = [
        ("3단계 적응 시퀀스",     "Phase 1 암흑 → Phase 2 엣지 → Phase 3 풀렌더링"),
        ("실험 조건 분기",        "대조군(1군) vs 실험군(2군) 자동 분기"),
        ("LobbyScene UI",         "XR 컨트롤러로 조건 선택, 씬 간 조건 유지"),
        ("피실험자 행동 통제",    "Phase별 안내 텍스트(InstructionUI) — 카메라 추적"),
        ("AimTrainer 콘텐츠",     "랜덤 타겟 스폰, 명중률/반응시간 측정"),
        ("CSV 데이터 로깅",       "Hit/Miss/Spawn 이벤트 — 조건명 자동 포함"),
    ]

    for i, (title, desc) in enumerate(items):
        top = Inches(1.42) + i * Inches(0.88)
        col = GREEN if i % 2 == 0 else BLUE
        pill(s, "✓ 완료", Inches(0.5), top + Inches(0.2),
             w=Inches(1.1), h=Inches(0.36), bg=col, size=12)
        box(s, title, Inches(1.75), top + Inches(0.06),
            Inches(3.0), Inches(0.42), size=15, bold=True, color=DARK)
        box(s, desc,  Inches(4.85), top + Inches(0.06),
            Inches(4.0), Inches(0.42), size=14, color=GRAY)

    # 우측 사진 2장
    safe_pic(s, img("로비.png"),
             Inches(9.1), Inches(1.38), width=Inches(3.95))
    safe_pic(s, img("대조군.png"),
             Inches(9.1), Inches(4.42), width=Inches(3.95), height=Inches(2.72))


# ════════════════════════════════════════════════════════════════
# 슬라이드 10: 실험 설계
# ════════════════════════════════════════════════════════════════
def s_exp(prs):
    s = blank(prs)
    header(s, "실험 설계", "집단 간 설계(Between-Subjects) · 피실험자 16명")

    # 1군 박스
    rect_bg(s, Inches(0.5), Inches(1.42), Inches(4.0), Inches(5.6), RGBColor(0xF5,0xF5,0xFF))
    pill(s, "1군 — 대조군", Inches(0.65), Inches(1.58),
         w=Inches(2.2), h=Inches(0.42), bg=BLUE)
    box(s, "HMD 착용\n→ 즉시 AimTrainer",
        Inches(0.65), Inches(2.15), Inches(3.7), Inches(0.9),
        size=20, bold=True, color=DARK)
    multibox(s, [
        "• 적응 절차 없음",
        "• fidelity = 1.0 고정",
        "• CSV: Control",
        "",
        "8명",
    ], Inches(0.65), Inches(3.1), Inches(3.7), Inches(3.5), size=17, color=DARK)

    box(s, "vs", Inches(4.6), Inches(3.5), Inches(0.8), Inches(0.8),
        size=24, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # 2군 박스
    rect_bg(s, Inches(5.5), Inches(1.42), Inches(4.0), Inches(5.6), RGBColor(0xF0,0xFF,0xF0))
    pill(s, "2군 — 실험군", Inches(5.65), Inches(1.58),
         w=Inches(2.2), h=Inches(0.42), bg=GREEN)
    box(s, "HMD 착용\n→ 3단계 적응 → AimTrainer",
        Inches(5.65), Inches(2.15), Inches(3.7), Inches(0.9),
        size=20, bold=True, color=DARK)
    multibox(s, [
        "• Phase 1 (30s) → Phase 2 (30s)",
        "• Phase 3 (30s) → AimTrainer",
        "• CSV: PostAdaptation",
        "",
        "8명",
    ], Inches(5.65), Inches(3.1), Inches(3.7), Inches(3.5), size=17, color=DARK)

    # 우측 — AimTrainer 사진
    safe_pic(s, img("대조군 조준훈련 1.png"),
             Inches(9.7), Inches(1.42), width=Inches(3.3))
    safe_pic(s, img("대조군 조준훈련2.png"),
             Inches(9.7), Inches(4.48), width=Inches(3.3), height=Inches(2.56))

    # 측정 도구 (하단)
    box(s, "측정:  VRSQ(멀미)  ·  SSQ(증상)  ·  SUS(실재감)  ·  CSV 행동 로그(명중률, 반응시간)",
        Inches(0.5), Inches(7.05), Inches(9.0), Inches(0.4),
        size=13, color=GRAY)


# ════════════════════════════════════════════════════════════════
# 슬라이드 11: 향후 계획
# ════════════════════════════════════════════════════════════════
def s_plan(prs):
    s = blank(prs)
    header(s, "앞으로 할 것들")

    phases_plan = [
        (PURPLE, "5월 4주", "Phase 1·2 몰입 시나리오 강화",
         "어둠 속 3D 공간음  /  윤곽선 뷰 행동 유도 시나리오"),
        (BLUE,   "5월 4주", "CQB 적 캐릭터 메시 교체",
         "구체 타겟 → 인체형 메시  /  Phase 2에서 실루엣으로 표현"),
        (ORANGE, "5월 4주", "SSQ / VRSQ 설문 시스템",
         "실험 전후 설문 자동 진행 및 CSV 병합"),
        (GREEN,  "6월 1주", "본 실험 진행  (16명)",
         "1군 8명 · 2군 8명  /  데이터 수집"),
        (DARK,   "6월 2주", "데이터 분석 및 논문 작성",
         "Independent t-test  /  상관 분석  /  고찰"),
    ]

    for i, (col, when, title, desc) in enumerate(phases_plan):
        top = Inches(1.38) + i * Inches(1.1)
        rect_bg(s, Inches(0.5), top, Inches(12.33), Inches(0.98),
                RGBColor(0xF8, 0xF8, 0xF8))
        pill(s, when, Inches(0.65), top + Inches(0.26),
             w=Inches(1.4), h=Inches(0.38), bg=col, size=12)
        box(s, title, Inches(2.25), top + Inches(0.1),
            Inches(4.5), Inches(0.45), size=17, bold=True, color=DARK)
        box(s, desc,  Inches(2.25), top + Inches(0.55),
            Inches(9.8), Inches(0.38), size=14, color=GRAY)


# ════════════════════════════════════════════════════════════════
# 슬라이드 12: 역할 분담
# ════════════════════════════════════════════════════════════════
def s_role(prs):
    s = blank(prs)
    header(s, "역할 분담",
           "공통: 연구 설계 · 선행 연구 · 실험 진행 · 데이터 분석 · 논문 작성 · 발표")

    members = [
        (BGPRP, PURPLE, "정근녕", "PM / Dev",
         ["Unity 아키텍처 설계", "적응 셰이더 · Phase 관리 시스템",
          "실험 분기 로직", "LobbyScene · InstructionUI", "시각 단계 전환 시스템"]),
        (BGBLUE, BLUE, "이수민", "Lead Dev",
         ["AimTrainer 콘텐츠 개발", "타겟 시스템 · HUD · 사격 로직",
          "데이터 로깅 파이프라인", "VR 사용자 조작 · 씬 구성", "향후: 공간 음향 시나리오"]),
        (BGGRN, GREEN, "곽경민", "UI / Env",
         ["CQB 맵 환경 디자인", "인터랙션 시스템 구현",
          "향후: 적 캐릭터 메시 구현", "피실험자 모집 · 일정 조율", "렌더링 연구"]),
    ]

    for i, (bg, col, name, role, items) in enumerate(members):
        lf = 0.5 + i * 4.25
        left = Inches(lf)
        rect_bg(s, left, Inches(1.42), Inches(3.9), Inches(5.65), bg)
        pill(s, role, left + Inches(0.18), Inches(1.58),
             w=Inches(1.5), h=Inches(0.4), bg=col, size=12)
        box(s, name, left + Inches(0.18), Inches(2.1),
            Inches(3.54), Inches(0.65), size=26, bold=True, color=DARK)
        line(s, Inches(2.85), left=lf + 0.2, right=lf + 3.7, color=col)
        for j, item in enumerate(items):
            box(s, f"• {item}",
                left + Inches(0.22), Inches(2.98) + j * Inches(0.74),
                Inches(3.46), Inches(0.65), size=14, color=DARK)


# ════════════════════════════════════════════════════════════════
# 슬라이드 13: 마무리
# ════════════════════════════════════════════════════════════════
def s_close(prs):
    s = blank(prs)
    box(s, "감사합니다",
        Inches(0.6), Inches(2.2), Inches(12), Inches(1.3),
        size=52, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    line(s, Inches(3.7))
    box(s, "VR 멀미 저감을 위한 교차 감각 기반 점진적 시각 동기화 기법",
        Inches(0.6), Inches(3.85), Inches(12), Inches(0.55),
        size=18, color=GRAY, align=PP_ALIGN.CENTER)
    box(s, "3팀  |  정근녕  ·  이수민  ·  곽경민",
        Inches(0.6), Inches(4.5), Inches(12), Inches(0.5),
        size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    box(s, "건국대학교 가상현실[3221]  |  2026년 5월 19일  |  담당: 김형석 교수님",
        Inches(0.6), Inches(5.2), Inches(12), Inches(0.45),
        size=14, color=LGRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 메인
# ════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    s_title(prs)     # 1
    s_agenda(prs)    # 2
    s_bg1(prs)       # 3
    s_bg2(prs)       # 4
    s_overview(prs)  # 5
    s_phase1(prs)    # 6
    s_phase2(prs)    # 7
    s_phase3(prs)    # 8
    s_impl(prs)      # 9
    s_exp(prs)       # 10
    s_plan(prs)      # 11
    s_role(prs)      # 12
    s_close(prs)     # 13

    out = r"D:\Git\SynchronizationInVR\Doc\중간_발표_PPT_초안.pptx"
    prs.save(out)
    print(f"saved: {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
